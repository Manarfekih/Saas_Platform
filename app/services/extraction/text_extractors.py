class TextExtractor:

    @staticmethod
    def extract_txt(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def extract_docx(file_path: str) -> str:
        try:
            from docx import Document as DocxDocument
        except ModuleNotFoundError as exc:
            raise RuntimeError(
                "DOCX extraction requires the optional 'python-docx' package."
            ) from exc

        doc = DocxDocument(file_path)
        text = []

        for paragraph in doc.paragraphs:
            stripped = paragraph.text.strip()
            if not stripped:
                continue

            style_name = (paragraph.style.name or "").lower()

            if "heading 1" in style_name or style_name == "title":
                text.append(f"## {stripped}")
            elif "heading" in style_name:
                text.append(f"### {stripped}")
            else:
                text.append(stripped)

        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(cells)

            if not rows:
                continue

            header = rows[0]
            text.append("| " + " | ".join(header) + " |")
            text.append("|" + "|".join(["---"] * len(header)) + "|")

            for row in rows[1:]:
                text.append("| " + " | ".join(row) + " |")

        return "\n".join(text)

    @staticmethod
    def extract_pptx(file_path: str) -> str:
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:
            raise RuntimeError(
                "PPTX extraction requires the optional 'python-pptx' package."
            ) from exc

        prs = Presentation(file_path)
        text_runs = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            text_runs.append(f"## Slide {slide_index}")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        text_runs.append(text)

        return "\n".join(text_runs)


text_extractor = TextExtractor()
