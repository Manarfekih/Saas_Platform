import os
import io
import json
import base64
from urllib import error, request

from PIL import Image
from pdf2image import convert_from_path

from docx import Document as DocxDocument
from pptx import Presentation


VISION_MODEL = os.getenv(
    "VISION_MODEL",
    "qwen2.5vl:7b"
)


OCR_DPI = int(os.getenv("OCR_DPI", "200"))

OCR_BATCH_SIZE = max(1, int(os.getenv("OCR_BATCH_SIZE", "1")))
OCR_MAX_SIDE = int(os.getenv("OCR_MAX_SIDE", "1500"))
OCR_REQUEST_TIMEOUT = int(os.getenv("OCR_REQUEST_TIMEOUT", "240"))

OLLAMA_URL = os.getenv(
    "OLLAMA_URL",
    "http://ollama:11434"
)



OCR_PROMPT = """
You are reading pages from the same document. The document type is
unknown and could be a resume, invoice, contract, report, letter,
form, or any other type. Transcribe all visible text exactly as it
appears. Do not summarize, paraphrase, or invent text.

LAYOUT RULES (critical):
- If the page has multiple columns, transcribe the ENTIRE left
  column first, top to bottom, then the ENTIRE right column, top to
  bottom. Never interleave lines from different columns.
- If the page has a header/footer separate from the main body,
  transcribe the main body first, then header/footer content at the
  very end of that page's output, clearly marked.
- Preserve reading order exactly as a human would read it: do not
  reorder sentences or merge unrelated blocks.

STRUCTURE RULES:
- When you see a section heading or label (e.g. "Experience",
  "Projects", "Education", "Skills", "Invoice Number", "Total Due",
  "Parties", "Clause 1"), output it as a markdown heading on its own
  line, like "## Experience", before its content.
- Preserve tables using markdown table syntax (| col | col |), keeping
  every row and column exactly as shown. Do not flatten tables into
  a single line of text.
- Preserve bullet points and numbered lists using "-" or "1." markers.

DATE-PAIRING RULES (critical):
- Only attach a date or date range to a title/line if that exact date
  is positioned directly adjacent to it in the image (same line,
  directly below it, or directly to its right/left as a clear visual
  pair). Do not infer or guess which title a date belongs to.
- If a date appears in its own separate row, block, or sub-column
  with no title directly touching it, transcribe it on its own
  separate line rather than merging it into the nearest heading or
  the previous item. It is better to leave a date unattached than to
  attach it to the wrong line.
- Never carry a date forward from one entry to a different entry. If
  several entries are listed with only one date visible, transcribe
  only the entry that the date is visually attached to with that
  date; leave the others without a date.

ITEM-PAIRING RULES (critical):
- Some sections contain SHORT side-by-side micro-blocks: a brief title
  or label (e.g. a course name, a certification name) positioned next
  to a separate, unrelated detail (e.g. a technology list or
  description belonging to a DIFFERENT item entirely). This is
  different from a simple two-column page layout — it can occur
  within a single section, as small adjacent pairs rather than full
  columns.
- Before pairing a title with a description/tech-list/bullet block,
  check whether they are genuinely the same item, or whether the
  title is short (a course/certification name with no further detail)
  while the description visually belongs to a different, adjacent
  item. If a short label (e.g. "DataCamp — Course Name") has no
  bullets or description directly beneath/after it before the next
  short label appears, transcribe it alone as a short standalone
  entry. Do not borrow the next visible description line to fill it
  in if that description is more plausibly attached to a different,
  larger item nearby.
- When uncertain whether two adjacent lines belong together, prefer
  transcribing them as separate standalone lines over merging them
  into one item. A human reader can still tell they're related from
  proximity; an incorrect merge cannot be undone downstream.

DIAGRAM RULES (critical):
- If the page contains a diagram, chart, flowchart, UML diagram,
  architecture diagram, org chart, or any other visual structure
  (boxes, arrows, connecting lines, nodes) rather than plain text:
  do NOT just list the text labels found inside it as a flat,
  unstructured sequence of words. A flat label dump destroys the
  diagram's meaning (e.g. listing actor names and action names from a
  UML use-case diagram with no indication of which actor performs
  which action).
- Instead, describe the diagram's structure in plain sentences: name
  the diagram type if identifiable (e.g. "Use case diagram", "Flow
  chart", "Entity relationship diagram"), then describe each
  connection or relationship explicitly, e.g. "Actor 'User' connects
  to use cases: Register, Login, Edit Profile, Upload Document, View
  Docs, Chat" rather than just "Register Login Edit Profile Upload
  Document View Docs Chat".
- If the diagram's relationships are too complex or unclear to
  describe with full confidence, still separate distinct labeled
  elements onto their own lines rather than concatenating them, and
  note "[DIAGRAM: structure unclear, labels listed individually]"
  before the list.
- Wrap diagram descriptions in their own subsection, e.g.
  "### Diagram: <short description of what it depicts>", so
  downstream processing can distinguish diagram content from regular
  prose.

OTHER RULES:
- If a page is blank, return only [BLANK].
- For each page, start with a marker like [[PAGE 1]] on its own line.
""".strip()




def _image_to_base64(img: Image.Image) -> str:

    if img.mode != "RGB":
        img = img.convert("RGB")

    buffer = io.BytesIO()
    img.save(
        buffer,
        format="JPEG",
        quality=85,
        optimize=True,
    )

    return base64.b64encode(
        buffer.getvalue()
    ).decode("utf-8")


def _prepare_image_for_ocr(img: Image.Image) -> Image.Image:

    prepared = img.copy()

    max_side = OCR_MAX_SIDE
    prepared.thumbnail((max_side, max_side), Image.Resampling.LANCZOS)

    if prepared.mode != "RGB":
        prepared = prepared.convert("RGB")

    return prepared


def _batch_items(items: list, batch_size: int):
    for index in range(0, len(items), batch_size):
        yield items[index:index + batch_size]


def _extract_from_images(
    images: list[Image.Image],
    page_start: int,
) -> str:

    image_b64s = [
        _image_to_base64(image)
        for image in images
    ]

    page_end = page_start + len(images) - 1
    page_markers = ", ".join(
        f"[[PAGE {page_num}]]"
        for page_num in range(page_start, page_end + 1)
    )

    payload = {
        "model": VISION_MODEL,
        "prompt": (
            f"{OCR_PROMPT}\n"
            f"The current images are pages {page_start} to {page_end}.\n"
            f"Use these markers in order: {page_markers}."
        ),
        "images": image_b64s,
        "stream": False,
        "options": {
            "temperature": 0,
            
            "num_ctx": 12288,
           
            "num_predict": 4096,
        },
    }

    req = request.Request(
        url=f"{OLLAMA_URL.rstrip('/')}/api/generate",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Content-Type": "application/json",
        },
        method="POST",
    )

    try:
        with request.urlopen(req, timeout=OCR_REQUEST_TIMEOUT) as resp:
            data = json.loads(resp.read().decode("utf-8"))

        text = data.get("response", "")

        
        done_reason = data.get("done_reason")

        if done_reason == "length":
            raise RuntimeError(
                f"OCR output truncated by length limit for pages "
                f"{page_start}-{page_end} (done_reason=length, "
                f"{len(text)} chars generated). Increase num_predict "
                f"or reduce OCR_BATCH_SIZE further."
            )

        return text.strip()

    except error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise RuntimeError(
            f"Ollama OCR request failed with HTTP {exc.code}: {detail}"
        ) from exc
    except error.URLError as exc:
        raise RuntimeError(
            f"Ollama OCR request failed: {exc.reason}"
        ) from exc


# normal text extractors

def _extract_txt(file_path: str) -> str:

    with open(
        file_path,
        "r",
        encoding="utf-8",
        errors="ignore",
    ) as f:
        return f.read()


def _extract_docx(file_path: str) -> str:
    

    doc = DocxDocument(file_path)

    text = []

    for paragraph in doc.paragraphs:
        stripped = paragraph.text.strip()

        if not stripped:
            continue

        style_name = (paragraph.style.name or "").lower()

        if "heading 1" in style_name or style_name == "title":
            text.append(f"## {stripped}")
        elif "heading" in style_name:
            text.append(f"### {stripped}")
        else:
            text.append(stripped)

    for table in doc.tables:

        rows = []

        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]

            if any(cells):
                rows.append(cells)

        if not rows:
            continue

        
        header = rows[0]
        text.append("| " + " | ".join(header) + " |")
        text.append("|" + "|".join(["---"] * len(header)) + "|")

        for row in rows[1:]:
            text.append("| " + " | ".join(row) + " |")

    return "\n".join(text)


def _extract_pptx(file_path: str) -> str:

    prs = Presentation(file_path)

    text_runs = []

    for slide_index, slide in enumerate(prs.slides, start=1):

        text_runs.append(f"## Slide {slide_index}")

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    text_runs.append(text)

    return "\n".join(text_runs)




def extract_text_llm(file_path: str) -> str:

    ext = file_path.lower().split(".")[-1]


    if ext == "txt":
        return _extract_txt(file_path)

    if ext == "docx":
        return _extract_docx(file_path)

    if ext == "pptx":
        return _extract_pptx(file_path)

    # Vision extraction only

    if ext == "pdf":

        images = convert_from_path(
            file_path,
            dpi=OCR_DPI,
            fmt="jpeg",
            grayscale=True,
            jpegopt={
                "quality": 85,
                "optimize": False,
                "progressive": False,
            },
            thread_count=max(1, OCR_BATCH_SIZE),
        )

    elif ext in ["png", "jpg", "jpeg"]:

        with Image.open(file_path) as img:
            images = [img.convert("RGB").copy()]

    else:
        raise ValueError(
            f"Unsupported format: {ext}"
        )

    if not images:
        return ""

    extracted_pages = []

    for batch_start_index, batch in enumerate(
        _batch_items(images, OCR_BATCH_SIZE),
        start=1,
    ):
        page_start = ((batch_start_index - 1) * OCR_BATCH_SIZE) + 1
        page_end = page_start + len(batch) - 1

        try:
            prepared_batch = [
                _prepare_image_for_ocr(image)
                for image in batch
            ]

            cleaned = _extract_from_images(
                prepared_batch,
                page_start=page_start,
            ).strip()

            if not cleaned:
                raise ValueError("Empty OCR output")

            extracted_pages.append(
                f"\n--- PAGES {page_start}-{page_end} ---\n{cleaned}"
            )

        except Exception as e:
            extracted_pages.append(
                f"\n--- PAGES {page_start}-{page_end} FAILED ---\n{str(e)}"
            )

    return "\n".join(extracted_pages)


def extract_text(file_path: str) -> str:
   

    return extract_text_llm(file_path)